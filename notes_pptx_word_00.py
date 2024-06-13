from pptx import Presentation
from docx import Document
import logging

# Initialize logging
logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')

def extract_and_write_notes_to_doc(pptx_file_path, docx_file_path):
    try:
        prs = Presentation(pptx_file_path)
        doc = Document()
    except Exception as e:
        logging.error(f"Error loading files: {e}")
        return

    for slide_number, slide in enumerate(prs.slides, start=1):
        try:
            notes_slide = slide.notes_slide
            if notes_slide and notes_slide.notes_text_frame:
                # Add slide number as Heading 1
                doc.add_heading(f'Slide {slide_number}', level=1)

                # Process each paragraph in the notes
                for paragraph in notes_slide.notes_text_frame.paragraphs:
                    # Create a new paragraph in the DOCX document for each paragraph in the notes
                    doc_paragraph = doc.add_paragraph()
                    # Add each run in the paragraph
                    for run in paragraph.runs:
                        doc_run = doc_paragraph.add_run(run.text)
                        # Copy formatting
                        doc_run.bold = run.font.bold
                        doc_run.italic = run.font.italic
                        # Here you could extend to other formatting like font size, etc.
            else:
                # Handle slides without notes
                logging.info(f"No notes found for Slide {slide_number}")

        except Exception as e:
            logging.warning(f"Error processing Slide {slide_number}: {e}")
            continue

    try:
        doc.save(docx_file_path)
        logging.info(f"Notes successfully written to {docx_file_path}")
    except Exception as e:
        logging.error(f"Error saving document: {e}")

pptx_file_path = 'Part 1 presentation_update00.pptx'
docx_file_path = 'presentation_update00_extracted_notes.docx'
extract_and_write_notes_to_doc(pptx_file_path, docx_file_path)
