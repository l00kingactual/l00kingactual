from pptx import Presentation

def extract_notes_from_presentation(pptx_file_path):
    # Load the presentation
    prs = Presentation(pptx_file_path)
    
    # Iterate through each slide
    for slide_number, slide in enumerate(prs.slides, start=1):
        # Access the notes slide for the current slide
        notes_slide = slide.notes_slide
        notes_text = notes_slide.notes_text_frame.text if notes_slide.notes_text_frame else "No notes for this slide."
        
        print(f"Slide {slide_number} Notes:")
        print(notes_text)
        print("-" * 40)

# Replace 'your_presentation.pptx' with the path to your PowerPoint file
pptx_file_path = 'your_presentation.pptx'
extract_notes_from_presentation(pptx_file_path)
