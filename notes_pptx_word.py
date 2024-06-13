from pptx import Presentation
from docx import Document
from docx.shared import Pt
import logging

# Initialize logging
logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')

def extract_and_write_notes_to_doc(pptx_file_path, docx_file_path):
    try:
        prs = Presentation(pptx_file_path)
        doc = Document()
    except Exception as e:
        logging.error(f"Error loading files: {e}")
        return

    for slide_number, slide in enumerate(prs.slides, start=1):
        try:
            notes_slide = slide.notes_slide
            notes_text = notes_slide.notes_text_frame.text if notes_slide.notes_text_frame else "No notes for this slide."

            # Add slide number as Heading 1
            doc.add_heading(f'Slide {slide_number}', level=1)

            # Add notes text as a paragraph, retaining simple formatting
            paragraph = doc.add_paragraph()
            for paragraph in notes_slide.notes_text_frame.paragraphs:
                run = paragraph.add_run(paragraph.text)
                # Copy basic formatting
                run.bold = paragraph.font.bold
                run.italic = paragraph.font.italic
                # Set font size if specified
                if paragraph.font.size:
                    run.font.size = Pt(paragraph.font.size.pt)

        except Exception as e:
            logging.warning(f"Error processing Slide {slide_number}: {e}")
            continue

    try:
        doc.save(docx_file_path)
        logging.info(f"Notes successfully written to {docx_file_path}")
    except Exception as e:
        logging.error(f"Error saving document: {e}")

pptx_file_path = 'Part 1 presentation_update00.pptx'
docx_file_path = 'presentation_update00_extracted_notes.docx'
extract_and_write_notes_to_doc(pptx_file_path, docx_file_path)
