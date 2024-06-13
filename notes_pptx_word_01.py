from docx import Document
from pptx import Presentation
import logging

# Setup logging
logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s: %(message)s')

def extract_notes_to_word(pptx_file, docx_file):
    try:
        prs = Presentation(pptx_file)
        doc = Document()
    except Exception as e:
        logging.error(f"Error loading presentation: {e}")
        return

    for slide_num, slide in enumerate(prs.slides, start=1):
        try:
            notes_slide = slide.notes_slide
            notes_text_frame = notes_slide.notes_text_frame

            # Add slide number as Heading
            doc.add_heading(f'Slide {slide_num}', level=1)
            
            if notes_text_frame:
                for paragraph in notes_text_frame.paragraphs:
                    # Create a new paragraph in the Word document for each paragraph in the notes
                    doc_para = doc.add_paragraph()
                    # Add each run (text segment) from the PPTX notes to the Word paragraph
                    for run in paragraph.runs:
                        doc_run = doc_para.add_run(run.text)
                        # Retain bold and italic formatting
                        doc_run.bold = run.font.bold
                        doc_run.italic = run.font.italic
            else:
                doc.add_paragraph("No notes for this slide.")
        except Exception as e:
            logging.warning(f"Error processing Slide {slide_num}: {e}")
    
    # Save the document
    try:
        doc.save(docx_file)
        logging.info(f"Notes successfully written to {docx_file}")
    except Exception as e:
        logging.error(f"Error saving the Word document: {e}")

pptx_file_path = 'Part 1 presentation_update00.pptx'
docx_file_path = 'presentation_update00_extracted_notes.docx'
extract_notes_to_word(pptx_file_path, docx_file_path)
